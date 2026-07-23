import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    DARK_BG = RGBColor(18, 24, 38)
    GOLD = RGBColor(212, 165, 116)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_GRAY = RGBColor(200, 210, 225)
    CARD_BG = RGBColor(28, 36, 56)

    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # SLIDE 1: Title
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide1, DARK_BG)

    top_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = GOLD
    top_bar.line.color.rgb = GOLD

    tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p0 = tf1.paragraphs[0]
    p0.text = "LIFEXTREME PREMIUM ADVENTURES"
    p0.font.bold = True
    p0.font.size = Pt(14)
    p0.font.color.rgb = GOLD

    p1 = tf1.add_paragraph()
    p1.text = "Atardecer en el Valle Sagrado"
    p1.font.bold = True
    p1.font.size = Pt(40)
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(10)

    p2 = tf1.add_paragraph()
    p2.text = "Moray & Salineras de Maras con Tea Time & Brindis Andino"
    p2.font.size = Pt(22)
    p2.font.color.rgb = GOLD
    p2.space_after = Pt(25)

    p3 = tf1.add_paragraph()
    p3.text = "Propuesta Comercial para Grupo de 50 Pasajeros | Fecha: 01 de Septiembre | Urubamba, Cusco"
    p3.font.size = Pt(14)
    p3.font.color.rgb = LIGHT_GRAY


    # SLIDE 2: Resumen
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide2, DARK_BG)

    tb_h2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf_h2 = tb_h2.text_frame
    p_h2 = tf_h2.paragraphs[0]
    p_h2.text = "Resumen de la Experiencia"
    p_h2.font.bold = True
    p_h2.font.size = Pt(28)
    p_h2.font.color.rgb = GOLD

    p_sub2 = tf_h2.add_paragraph()
    p_sub2.text = "Diseñado para el día de llegada: Aclimatación suave, relajación y paisajes impresionantes"
    p_sub2.font.size = Pt(14)
    p_sub2.font.color.rgb = LIGHT_GRAY

    cards_data = [
        {"title": "🌸 Aclimatación Suave", "desc": "Caminatas ligeras a tu propio ritmo. Ideal para adaptarse a la altitud del Valle Sagrado sin esfuerzo físico."},
        {"title": "📸 Paisajes Increíbles", "desc": "Visita a las terrazas circulares de Moray y a las milenarias pozas de sal rosada de Maras durante la hora dorada."},
        {"title": "☕ Experiencia Tea Time", "desc": "Pausa reconfortante en mirador privado con infusiones andinas digestivas, snacks locales y brindis de bienvenida."}
    ]

    for i, c in enumerate(cards_data):
        left = Inches(0.8 + i * 3.9)
        top = Inches(1.8)
        width = Inches(3.6)
        height = Inches(4.8)

        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = GOLD

        tb = slide2.shapes.add_textbox(left + Inches(0.2), top + Inches(0.3), width - Inches(0.4), height - Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True

        p_t = tf.paragraphs[0]
        p_t.text = c["title"]
        p_t.font.bold = True
        p_t.font.size = Pt(18)
        p_t.font.color.rgb = GOLD
        p_t.space_after = Pt(14)

        p_d = tf.add_paragraph()
        p_d.text = c["desc"]
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = WHITE


    # SLIDE 3: Itinerario
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide3, DARK_BG)

    tb_h3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf_h3 = tb_h3.text_frame
    p_h3 = tf_h3.paragraphs[0]
    p_h3.text = "Itinerario Detallado (14:00 hrs – 19:00 hrs)"
    p_h3.font.bold = True
    p_h3.font.size = Pt(28)
    p_h3.font.color.rgb = GOLD

    timeline = [
        ("14:00 hrs", "Recojo en Hotel (Urubamba)", "Recepción en 2 buses Coaster exclusivos con atención personalizada."),
        ("14:30 hrs", "Laboratorio Inca de Moray", "Recorrido guiado relajado por los impresionantes andenes circulares."),
        ("16:00 hrs", "Salineras de Maras", "Vista panorámica de las 3,000+ pozas de sal rosada y tiempo de fotos."),
        ("17:15 hrs", "Tea Time & Brindis Andino", "Pausa en mirador privado con mates de muña/coca, snacks y copa de bienvenida."),
        ("18:30 - 19:00 hrs", "Retorno al Hotel", "Viaje con vistas al atardecer sobre los nevados y llegada a Urubamba.")
    ]

    for i, (time_str, title_str, desc_str) in enumerate(timeline):
        top = Inches(1.6 + i * 1.05)
        
        time_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(2.2), Inches(0.8))
        time_box.fill.solid()
        time_box.fill.fore_color.rgb = GOLD
        time_box.line.color.rgb = GOLD
        tf_time = time_box.text_frame
        p_time = tf_time.paragraphs[0]
        p_time.text = time_str
        p_time.font.bold = True
        p_time.font.size = Pt(14)
        p_time.font.color.rgb = DARK_BG
        p_time.alignment = PP_ALIGN.CENTER

        content_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), top, Inches(9.3), Inches(0.8))
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = CARD_BG
        content_box.line.color.rgb = GOLD
        tf_content = content_box.text_frame
        tf_content.word_wrap = True

        p_ct = tf_content.paragraphs[0]
        p_ct.text = title_str + " — "
        p_ct.font.bold = True
        p_ct.font.size = Pt(15)
        p_ct.font.color.rgb = GOLD

        run_desc = p_ct.add_run()
        run_desc.text = desc_str
        run_desc.font.bold = False
        run_desc.font.size = Pt(13)
        run_desc.font.color.rgb = WHITE


    # SLIDE 4: Servicios Incluidos
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide4, DARK_BG)

    tb_h4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf_h4 = tb_h4.text_frame
    p_h4 = tf_h4.paragraphs[0]
    p_h4.text = "Servicios Incluidos & Garantía Lifextreme"
    p_h4.font.bold = True
    p_h4.font.size = Pt(28)
    p_h4.font.color.rgb = GOLD

    inc_left = [
        "✔ 2 Buses Turísticos Coaster exclusivos (AC + Balón de oxígeno)",
        "✔ 2 Guías Oficiales Bilingües (1 guía dedicado por bus)",
        "✔ Boletos de ingreso completo a Moray y Salineras de Maras"
    ]
    inc_right = [
        "✔ Experiencia Tea Time: Mates digestivos + Snacks artesanales",
        "✔ Brindis de Bienvenida al atardecer en mirador privado",
        "✔ Asistencia y coordinación logística en tiempo real"
    ]

    box_l = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = CARD_BG
    box_l.line.color.rgb = GOLD
    tf_l = box_l.text_frame
    tf_l.word_wrap = True
    
    p_l_t = tf_l.paragraphs[0]
    p_l_t.text = "Logística & Guiado (2 Buses)"
    p_l_t.font.bold = True
    p_l_t.font.size = Pt(18)
    p_l_t.font.color.rgb = GOLD
    p_l_t.space_after = Pt(14)

    for item in inc_left:
        p_item = tf_l.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = WHITE
        p_item.space_after = Pt(12)

    box_r = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = CARD_BG
    box_r.line.color.rgb = GOLD
    tf_r = box_r.text_frame
    tf_r.word_wrap = True

    p_r_t = tf_r.paragraphs[0]
    p_r_t.text = "Experiencia & Entradas"
    p_r_t.font.bold = True
    p_r_t.font.size = Pt(18)
    p_r_t.font.color.rgb = GOLD
    p_r_t.space_after = Pt(14)

    for item in inc_right:
        p_item = tf_r.add_paragraph()
        p_item.text = item
        p_item.font.size = Pt(14)
        p_item.font.color.rgb = WHITE
        p_item.space_after = Pt(12)


    # SLIDE 5: Presupuesto
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(slide5, DARK_BG)

    tb_h5 = slide5.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.0))
    tf_h5 = tb_h5.text_frame
    p_h5 = tf_h5.paragraphs[0]
    p_h5.text = "Inversión Comercial (Incluye IGV 18%)"
    p_h5.font.bold = True
    p_h5.font.size = Pt(28)
    p_h5.font.color.rgb = GOLD

    price_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.3))
    price_card.fill.solid()
    price_card.fill.fore_color.rgb = CARD_BG
    price_card.line.color.rgb = GOLD

    tf_pc = price_card.text_frame
    tf_pc.word_wrap = True

    p_pc1 = tf_pc.paragraphs[0]
    p_pc1.text = "PRECIO POR PERSONA:  USD $65.00 + IGV ($11.70) = USD $76.70"
    p_pc1.font.bold = True
    p_pc1.font.size = Pt(20)
    p_pc1.font.color.rgb = GOLD
    p_pc1.alignment = PP_ALIGN.CENTER
    p_pc1.space_after = Pt(10)

    p_pc2 = tf_pc.add_paragraph()
    p_pc2.text = "INVERSIÓN TOTAL DEL GRUPO (50 PAX INC. IGV):  USD $3,835.00"
    p_pc2.font.bold = True
    p_pc2.font.size = Pt(26)
    p_pc2.font.color.rgb = WHITE
    p_pc2.alignment = PP_ALIGN.CENTER

    cond_card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.2))
    cond_card.fill.solid()
    cond_card.fill.fore_color.rgb = CARD_BG
    cond_card.line.color.rgb = GOLD

    tf_cc = cond_card.text_frame
    tf_cc.word_wrap = True

    p_cc_t = tf_cc.paragraphs[0]
    p_cc_t.text = "Condiciones de Reserva & Emisión de Comprobante"
    p_cc_t.font.bold = True
    p_cc_t.font.size = Pt(16)
    p_cc_t.font.color.rgb = GOLD
    p_cc_t.space_after = Pt(8)

    conds = [
        "• Emisión de Factura o Boleta de Venta electrónica con IGV (18%) incluido.",
        "• Reserva con el 50% de adelanto y confirmación de nómina final 7 días antes de la llegada.",
        "• Cancelación sin penalidad hasta 10 días antes del inicio del servicio."
    ]
    for cond in conds:
        p_c = tf_cc.add_paragraph()
        p_c.text = cond
        p_c.font.size = Pt(13)
        p_c.font.color.rgb = LIGHT_GRAY

    output_path = r"C:\Users\ASUS\OneDrive\VARIOS\Documentos\GPTS IA\BIOVET AI\Lifextreme-Web-AI\Propuesta_Valle_Sagrado_50pax.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_deck()
